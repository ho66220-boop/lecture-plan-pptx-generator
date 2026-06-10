# -*- coding: utf-8 -*-
"""강사 사진을 우상단 사진 박스(roundRect)에 채워 넣는다.

- 사진은 `teacher_photo_dir` 폴더에 강사명으로 저장(예: 홍길동.jpg).
- 박스의 solidFill을 picture fill(blipFill)로 교체 → 둥근 모서리·테두리를 유지한 채 사진이 들어감.
- 정사각 박스에 맞춰 가운데 기준으로 잘라(cover) 비율 왜곡 없이 채운다.
- 사진이 없으면 아무것도 하지 않는다(회색 박스 유지).
"""
import re
import unicodedata
from pathlib import Path

from pptx.oxml.ns import qn

EMU_PER_CM = 360000.0
PHOTO_EXTS = {".jpg", ".jpeg", ".png", ".webp"}

# 우상단 정사각 사진 박스 식별 기준(reflow·photo 공용 — 한 군데에서만 정의).
PHOTO_BOX_MAX_TOP = 2.0 * EMU_PER_CM     # 이보다 위(상단)
PHOTO_BOX_MIN_LEFT = 13.5 * EMU_PER_CM   # 이보다 오른쪽
PHOTO_BOX_MIN_SIZE = 3.0 * EMU_PER_CM    # 최소 한 변 길이
PHOTO_BOX_SQUARE_TOL = 0.5 * EMU_PER_CM  # 가로·세로 차이 허용(정사각 판정)


def is_photo_box(shape):
    """우상단의 빈(텍스트 없는) 정사각형 사진 박스인지 위치·형태로 판별.
    teacher_photo(사진 채우기 대상 탐색)와 reflow(높이 확장 보호) 양쪽에서 공유한다."""
    if shape.top is None or shape.left is None or shape.width is None or shape.height is None:
        return False
    if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
        return False
    return (
        shape.top < PHOTO_BOX_MAX_TOP
        and shape.left > PHOTO_BOX_MIN_LEFT
        and shape.width > PHOTO_BOX_MIN_SIZE
        and abs(shape.width - shape.height) < PHOTO_BOX_SQUARE_TOL
    )


def _name_key(stem):
    """사진 파일명에서 강사 '이름'만 추출해 정규화.

    접두사는 무엇이든 무관(마지막 '_' 뒤를 이름으로 본다). 지원 형식:
      홍길동.jpg                    -> 홍길동
      academy_국어_홍길동2.png       -> 홍길동   (배포처_과목_이름 + 뒤 숫자)
      academy_수리논술_홍길동2 (1).png -> 홍길동
      academy_수학_홍길동[1].png     -> 홍길동
    """
    s = unicodedata.normalize("NFC", stem)
    if "_" in s:
        s = s.rsplit("_", 1)[-1]            # 배포처_과목_이름 -> 이름...
    s = re.sub(r"[\(\[][^)\]]*[\)\]]", "", s)  # (1), [1], (2) 등 중복표시 제거
    s = re.sub(r"\d+$", "", s.strip())          # 끝에 붙은 숫자 제거(홍길동2 -> 홍길동)
    return s.strip()

# spPr 안에서 fill 요소 뒤에 와야 하는 자식들(이 중 첫 요소 앞에 blipFill을 끼워 넣는다)
_POST_FILL_TAGS = ("a:ln", "a:effectLst", "a:effectDag", "a:scene3d", "a:sp3d", "a:extLst")
_FILL_TAGS = ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill")


def find_photo_box(slide):
    """우상단의 정사각형 사진 박스를 찾는다(여럿이면 가장 큰 것)."""
    best = None
    for sh in slide.shapes:
        if is_photo_box(sh):
            if best is None or (sh.width * sh.height) > (best.width * best.height):
                best = sh
    return best


def find_teacher_photo(photo_dir, teacher_name):
    """teacher_photo_dir에서 강사명과 일치하는 사진 경로를 찾는다(NFC 정규화 비교)."""
    if not photo_dir or not teacher_name:
        return None
    directory = Path(photo_dir)
    if not directory.is_dir():
        return None
    target = _name_key(str(teacher_name))
    if not target:
        return None
    for file_path in sorted(directory.iterdir()):
        if file_path.suffix.lower() in PHOTO_EXTS and _name_key(file_path.stem) == target:
            return str(file_path)
    return None


def _cover_src_rect(box_shape, image):
    """박스 비율에 맞춰 가운데 crop(cover)할 srcRect 속성(per-mille)을 계산."""
    iw, ih = image.size
    if not iw or not ih:
        return {}
    src_aspect = iw / ih
    box_aspect = box_shape.width / box_shape.height
    left = top = right = bottom = 0.0
    if src_aspect > box_aspect:        # 이미지가 더 넓음 → 좌우 crop
        crop = (1 - box_aspect / src_aspect) / 2
        left = right = crop
    elif src_aspect < box_aspect:      # 이미지가 더 김 → 상하 crop
        crop = (1 - src_aspect / box_aspect) / 2
        top = bottom = crop
    attrs = {}
    for key, val in (("l", left), ("t", top), ("r", right), ("b", bottom)):
        if val > 0:
            attrs[key] = str(int(round(val * 100000)))
    return attrs


def set_picture_fill(box_shape, slide, image_path):
    """박스의 채우기를 사진으로 교체(둥근 모서리·테두리 유지, 가운데 crop)."""
    image_part, rId = slide.part.get_or_add_image_part(image_path)
    sp_pr = box_shape._element.spPr

    for tag in _FILL_TAGS:
        for elem in sp_pr.findall(qn(tag)):
            sp_pr.remove(elem)

    blip_fill = sp_pr.makeelement(qn("a:blipFill"), {})
    blip = sp_pr.makeelement(qn("a:blip"), {qn("r:embed"): rId})
    blip_fill.append(blip)

    src_attrs = _cover_src_rect(box_shape, image_part.image)
    if src_attrs:
        blip_fill.append(sp_pr.makeelement(qn("a:srcRect"), src_attrs))

    stretch = sp_pr.makeelement(qn("a:stretch"), {})
    stretch.append(sp_pr.makeelement(qn("a:fillRect"), {}))
    blip_fill.append(stretch)

    anchor = None
    for tag in _POST_FILL_TAGS:
        found = sp_pr.find(qn(tag))
        if found is not None:
            anchor = found
            break
    if anchor is not None:
        anchor.addprevious(blip_fill)
    else:
        sp_pr.append(blip_fill)


def apply_teacher_photo(slide, teacher_name, photo_dir):
    """강사 사진을 사진 박스에 넣는다. 넣었으면 사진 경로, 아니면 None 반환."""
    if not photo_dir:
        return None
    photo_path = find_teacher_photo(photo_dir, teacher_name)
    if not photo_path:
        return None
    box = find_photo_box(slide)
    if box is None:
        return None
    set_picture_fill(box, slide, photo_path)
    return photo_path
