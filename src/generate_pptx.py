import copy
import re
from datetime import date, datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

try:
    from src.date_utils import format_date_dot
    from src.reflow import reflow_slide
    from src.subject_band import add_subject_band
    from src.teacher_photo import apply_teacher_photo
    from src.validate import report_row
except ModuleNotFoundError:
    from .date_utils import format_date_dot
    from .reflow import reflow_slide
    from .subject_band import add_subject_band
    from .teacher_photo import apply_teacher_photo
    from .validate import report_row


PLACEHOLDER_RE = re.compile(r"\{\{([^{}]+)\}\}")
DEFAULT_TEMPLATE_PATH = Path("templates") / "강의계획서_마스터템플릿.pptx"

# ── A4 맞춤(슬라이드는 프레젠테이션당 한 크기 → A4 고정, 넘치면 본문 폰트 축소) ──
EMU_PER_CM = 360000
EMU_PER_PT = 12700
A4_BOTTOM_MARGIN = int(0.30 * EMU_PER_CM)   # A4 하단 여유
FIT_TITLE_GUARD = int(5.0 * EMU_PER_CM)     # 이 위(제목·과목 영역)는 축소 대상에서 제외
FIT_SCALES = [1.0, 0.95, 0.90, 0.85, 0.80, 0.75, 0.70, 0.65, 0.60, 0.55, 0.50]
PROSE_MIN_WIDTH = int(15.0 * EMU_PER_CM)    # 이보다 넓은 글상자 = 강의특징/관리 같은 본문 → 먼저 축소
PROSE_MIN_CHARS = 40                        # 본문으로 볼 최소 글자 수(1단어 헤더는 제외)
REST_FLOOR = 0.82                           # 정보 그리드·진도표는 이 비율 아래로는 줄이지 않음(가독성 유지)
PROSE_MIN_PT = 11.0                         # 강의특징/관리 본문 폰트 절대 하한(pt). base가 더 작으면 base에서 멈춤.
                                            # 하한으로도 한 장에 안 들어가면 자르지 않고 A4_OVERFLOW로 플래그만.

# ── 좌상단 배지(강좌유형/구분/학년)·과목/강사 한 줄 박스: 가로로 넘치면 폰트만 축소 ──
BADGE_FIT_TOP = FIT_TITLE_GUARD             # 이 위(상단)의 좁은 한 줄 박스만 가로 맞춤 대상
BADGE_MAX_WIDTH = int(4.5 * EMU_PER_CM)     # 이보다 넓은 박스(슬로건/제목)는 제외
BADGE_FONT_FLOOR = 7.0                      # 이 pt 아래로는 줄이지 않음
DEFAULT_INSET = 91440                       # 좌/우 기본 안쪽 여백(EMU)
# 큰 제목(강좌명): 박스(높이 2.5cm)가 2줄을 수용하므로, 2줄 안에 들면 폰트 유지하고
# 그보다 길 때만 축소(긴 강좌명도 큰 글씨로 최대 2줄까지).
TITLE_MIN_WIDTH = int(10.0 * EMU_PER_CM)    # 이보다 넓은 박스만 제목으로 간주(슬로건/제목)
TITLE_BIG_PT = 16.0                         # 이 pt 이상이어야 '큰 제목'(작은 슬로건 글상자는 제외)
TITLE_FONT_FLOOR = 14.0                     # 큰 제목 최소 pt
TITLE_MAX_LINES = 2                         # 제목 허용 최대 줄 수

# 수학 세부 과목 → '수학' 으로 통일할 때 쓰는 키워드(공백 제거 후 비교).
MATH_SUBJECT_KEYS = (
    "수학", "수1", "수2", "수상", "수하", "미적", "확통", "확률",
    "통계", "기하", "대수", "공통수학",
)
# 끝에 붙은 수준 표기 제거: 아라비아(1·2·3)·전각(１２３)·로마자(Ⅰ Ⅱ Ⅲ).
# 예) 물리학Ⅰ→물리학, 화학Ⅰ→화학, 통합과학2→통합과학, 생명과학Ⅰ→생명과학.
LEVEL_SUFFIX_RE = re.compile(r"\s*[123１２３ⅠⅡⅢ]\s*$")


def normalize_subject(subject):
    """과목 표기 정규화:
    - 수학 계열 세부 과목(미적분·확통·기하 등)은 모두 '수학' 으로.
    - 과탐 등 끝에 붙은 수준 표기(물리학Ⅰ·화학Ⅰ·통합과학2 …) 제거.
    """
    s = (subject or "").strip()
    if not s:
        return s
    key = s.replace(" ", "")
    if any(word in key for word in MATH_SUBJECT_KEYS):
        return "수학"
    return LEVEL_SUFFIX_RE.sub("", s)


def _text_visual_pt(text, font_pt):
    """한 줄 텍스트의 대략적 가로폭(pt). 한글=폰트pt, ASCII=절반 가량."""
    return sum((font_pt * 0.55) if ord(ch) < 128 else float(font_pt) for ch in text)


def fit_oneline_badges(slide):
    """상단의 좁은 한 줄 박스(배지·과목·강사명)에서 글자가 박스 폭을 넘치면
    폰트만 줄여 한 줄에 들어가게 한다(세로 확장/줄바꿈 방지). 리플로우 전에 호출."""
    for shape in slide.shapes:
        if shape.top is None or shape.top >= BADGE_FIT_TOP:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.width is None or shape.width > BADGE_MAX_WIDTH:
            continue
        tf = shape.text_frame
        text = tf.text.strip()
        if not text:
            continue
        tf.word_wrap = False
        runs = [run for para in tf.paragraphs for run in para.runs]
        if not runs:
            continue
        base = next((r.font.size.pt for r in runs if r.font.size is not None), 9.0)
        ml = tf.margin_left if tf.margin_left is not None else DEFAULT_INSET
        mr = tf.margin_right if tf.margin_right is not None else DEFAULT_INSET
        inner_pt = (shape.width - ml - mr) / EMU_PER_PT
        widest = max((_text_visual_pt(line, base) for line in text.split("\n")), default=0)
        if inner_pt <= 0 or widest <= inner_pt:
            continue
        new_pt = max(BADGE_FONT_FLOOR, round(base * inner_pt / widest, 1))
        for run in runs:
            if run.font.size is not None:
                run.font.size = Pt(new_pt)


def fit_title(slide):
    """큰 제목(강좌명) 글상자: 박스 높이(2.5cm)가 최대 2줄을 담으므로, 2줄 안에 들어가면
    폰트를 그대로 두고(큰 글씨 유지), 2줄로도 부족할 만큼 길면 그때만 폰트를 줄인다."""
    for shape in slide.shapes:
        if shape.top is None or shape.top >= BADGE_FIT_TOP:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.width is None or shape.width < TITLE_MIN_WIDTH:
            continue
        tf = shape.text_frame
        paragraphs = [para.text for para in tf.paragraphs if para.text.strip()]
        if not paragraphs:
            continue
        runs = [run for para in tf.paragraphs for run in para.runs]
        base = next((r.font.size.pt for r in runs if r.font.size is not None), 0.0)
        if base < TITLE_BIG_PT:        # 작은 슬로건 글상자는 제외
            continue
        tf.word_wrap = True            # 2줄 줄바꿈 허용
        ml = tf.margin_left if tf.margin_left is not None else DEFAULT_INSET
        mr = tf.margin_right if tf.margin_right is not None else DEFAULT_INSET
        inner_pt = (shape.width - ml - mr) / EMU_PER_PT
        widest = max((_text_visual_pt(p, base) for p in paragraphs), default=0)
        if inner_pt <= 0 or widest <= 0:
            continue
        # 최대 줄 수 안에 들어가는 최대 폰트. 이미 들어가면(>=base) 손대지 않음.
        max_font = base * (TITLE_MAX_LINES * inner_pt) / widest
        if max_font >= base:
            continue
        new_pt = max(TITLE_FONT_FLOOR, round(max_font, 1))
        for run in runs:
            if run.font.size is not None:
                run.font.size = Pt(new_pt)


def _is_prose_shape(shape):
    return (
        getattr(shape, "has_text_frame", False)
        and shape.width is not None
        and shape.width > PROSE_MIN_WIDTH
        and len(shape.text_frame.text.strip()) > PROSE_MIN_CHARS
    )


def _snapshot_layout(slide):
    """리플로우 전 상태(도형 위치·높이, 본문 런 폰트pt)를 저장. 폰트 축소 재시도용."""
    pos, fonts = {}, {}
    for shape in slide.shapes:
        if shape.top is not None and shape.height is not None:
            pos[shape.shape_id] = (shape.top, shape.height)
        if getattr(shape, "has_text_frame", False):
            for pi, para in enumerate(shape.text_frame.paragraphs):
                for ri, run in enumerate(para.runs):
                    if run.font.size is not None:
                        fonts[(shape.shape_id, pi, ri)] = run.font.size.pt
    return pos, fonts


def _restore_layout(slide, snapshot):
    pos, fonts = snapshot
    for shape in slide.shapes:
        if shape.shape_id in pos:
            shape.top, shape.height = pos[shape.shape_id]
        if getattr(shape, "has_text_frame", False):
            for pi, para in enumerate(shape.text_frame.paragraphs):
                for ri, run in enumerate(para.runs):
                    key = (shape.shape_id, pi, ri)
                    if key in fonts:
                        run.font.size = Pt(fonts[key])


def _max_bottom(slide):
    bottom = 0
    for shape in slide.shapes:
        if shape.top is not None and shape.height is not None:
            bottom = max(bottom, shape.top + shape.height)
    return bottom


def fit_slide_to_height(slide, content_ids, target_bottom, guard=FIT_TITLE_GUARD):
    """내용이 target_bottom(A4) 안에 들도록 본문 폰트를 단계적으로 축소하며 리플로우.
    제목·과목 영역(guard 위)은 축소하지 않는다. 반환: (적용 scale, 맞췄는지)."""
    snapshot = _snapshot_layout(slide)
    fonts = snapshot[1]
    last_scale = 1.0
    for scale in FIT_SCALES:
        _restore_layout(slide, snapshot)
        if scale < 1.0:
            # 본문(강의특징/관리)은 절대 하한 PROSE_MIN_PT까지만, 정보·진도표는 REST_FLOOR 비율까지만 줄여 가독성 유지.
            rest_scale = max(scale, REST_FLOOR)
            for shape in slide.shapes:
                if shape.top is None or shape.top < guard:
                    continue
                if not getattr(shape, "has_text_frame", False):
                    continue
                prose = _is_prose_shape(shape)
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    for ri, run in enumerate(para.runs):
                        key = (shape.shape_id, pi, ri)
                        if key not in fonts:
                            continue
                        base_pt = fonts[key]
                        if prose:
                            floor_pt = min(base_pt, PROSE_MIN_PT)   # base가 floor보다 작으면 base에서 멈춤
                            new_pt = max(floor_pt, base_pt * scale)
                        else:
                            new_pt = base_pt * rest_scale
                        run.font.size = Pt(round(new_pt, 1))
        reflow_slide(slide, content_ids)
        last_scale = scale
        if _max_bottom(slide) <= target_bottom:
            return scale, True
    return last_scale, False


def clone_template_slide(prs, template_slide):
    """Clone shapes from the template slide into a new blank slide."""
    blank_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape.element)
    for shape in template_slide.shapes:
        slide.shapes._spTree.insert_element_before(copy.deepcopy(shape.element), "p:extLst")
    return slide


def remove_slide(prs, slide):
    slide_id = slide.slide_id
    slide_element = prs.slides._sldIdLst
    rel_id = None
    for item in slide_element:
        if int(item.id) == slide_id:
            rel_id = item.rId
            slide_element.remove(item)
            break
    if rel_id:
        prs.part.drop_rel(rel_id)


def text_value(value):
    if value is None:
        return ""
    return str(value)


def progress_date_display(row):
    # normalize 단계에서 base_year를 적용해 계산해 둔 parsed_date(ISO)를 그대로 사용한다.
    # (여기서 parse_date를 다시 호출하면 base_year 기본값으로 재파싱되어 요일이 틀어질 수 있음)
    iso = row.get("parsed_date")
    if iso:
        return format_date_dot(date.fromisoformat(iso))
    return text_value(row.get("날짜", ""))


HOLIDAY_RED = RGBColor(0xC0, 0x00, 0x00)   # 진도계획 휴강 표시 색


def progress_content_display(row):
    values = [row.get("수업 주제", ""), row.get("상세 내용", ""), row.get("비고", "")]
    lines = []
    for value in values:
        text = text_value(value).strip()
        if text and text not in lines:   # 중복 줄 제거(휴강\n휴강 → 휴강)
            lines.append(text)
    return "\n".join(lines)


PROGRESS_MAX_PER_COL = 5     # 한 컬럼(좌/우) 최대 행 수
PROGRESS_ROW_GAP_CM = 0.45   # 진도표 행 클러스터링 임계(행 내 도형 간격 < 이 값 < 행 사이 간격)


def progress_split(n):
    """진도 n개를 좌우 균형 분배. 좌측=ceil(n/2)(최대 5), 우측=나머지(최대 5).
    읽기 순서: 좌측 위→아래가 앞 회차, 그다음 우측. 예) 4→(2,2), 9→(5,4), 10→(5,5)."""
    n = min(n, PROGRESS_MAX_PER_COL * 2)
    left = min(PROGRESS_MAX_PER_COL, (n + 1) // 2)
    return left, n - left


def _cluster_by_top(shapes, gap_emu):
    """top 기준으로 도형을 '행'으로 묶는다(연속 정렬 후 간격 > gap이면 새 행)."""
    rows = []
    for sh in sorted(shapes, key=lambda s: s.top):
        if rows and sh.top - rows[-1][-1].top <= gap_emu:
            rows[-1].append(sh)
        else:
            rows.append([sh])
    return rows


def balance_progress_columns(slide, left_count, right_count):
    """진도표에서 채우지 않는 칸의 도형(배경+글상자)을 삭제한다. 칸 식별은 id가 아니라 위치:
    진도 섹션 제목 아래 도형을 top으로 묶어 '행', left로 갈라 '좌/우 컬럼'을 판별.
    우측이 한 칸도 없을 때(right_count==0)만 우측 헤더도 삭제."""
    title = None
    for sh in slide.shapes:                              # 진도 섹션 제목(가장 아래의 '진도' 텍스트)
        if getattr(sh, "has_text_frame", False) and sh.top is not None and "진도" in sh.text_frame.text:
            if title is None or sh.top > title.top:
                title = sh
    if title is None:
        return
    min_dim = int(0.3 * EMU_PER_CM)
    region = [
        sh for sh in slide.shapes
        if sh.top is not None and sh.height is not None and sh.width is not None
        and sh.top > title.top + min_dim          # 제목보다 아래
        and sh.width > min_dim and sh.height > min_dim   # 장식 바(얇은 선) 제외
    ]
    rows = _cluster_by_top(region, int(PROGRESS_ROW_GAP_CM * EMU_PER_CM))
    if len(rows) < 2:
        return
    header, data_rows = rows[0], rows[1 : 1 + PROGRESS_MAX_PER_COL]

    lefts = sorted({sh.left for sh in region})           # 좌/우 컬럼 경계 = 가장 큰 left 간격
    boundary, best = None, -1
    for a, b in zip(lefts, lefts[1:]):
        if b - a > best:
            best, boundary = b - a, (a + b) // 2
    if boundary is None:
        return

    to_delete = []
    for r, row in enumerate(data_rows, start=1):
        if r > left_count:
            to_delete += [sh for sh in row if sh.left < boundary]
        if r > right_count:
            to_delete += [sh for sh in row if sh.left >= boundary]
    if right_count == 0:                                  # 우측 통째로 안 쓰면 우측 헤더도 제거
        to_delete += [sh for sh in header if sh.left >= boundary]
    for sh in to_delete:
        sh._element.getparent().remove(sh._element)


def build_placeholder_map(lecture):
    fields = lecture.get("fields", {})
    # 큰 글씨 = 강좌명, 작은 글씨 = 제목 우선(단 제목이 비었거나 강좌명과 같으면 서브 슬로건).
    course_name = fields.get("강좌명", "")
    main_title = fields.get("메인 제목", "")
    sub_slogan = fields.get("서브 슬로건", "")
    small_caption = (
        main_title
        if main_title.strip() and main_title.strip() != course_name.strip()
        else sub_slogan
    )
    placeholder_map = {
        "강좌유형": fields.get("강좌 유형", ""),
        "구분": fields.get("구분", ""),
        "학년": fields.get("학년", ""),
        "강의형태": fields.get("강의형태", ""),
        "서브슬로건": small_caption,   # 작은 글씨(상단)
        "메인제목": course_name,       # 큰 글씨 = 강좌명
        "과목": normalize_subject(fields.get("과목", "")),
        "강사명": fields.get("강사명", ""),
        "강좌명": fields.get("강좌명", ""),
        "개강일": lecture.get("opening_date_display", ""),
        "수업시간": fields.get("수업 요일 / 시간", ""),
        "수강기간": lecture.get("period_display", ""),
        "휴강일": fields.get("휴강일 / 사유 / 수업 불가 일정", ""),
        "보강방법": fields.get("보강 방법", ""),
        "수강료": lecture.get("fee_display", ""),
        "교재정보": fields.get("교재 정보", ""),
        "강의특징": fields.get("강의특징", ""),
        "관리프로그램": fields.get("관리 프로그램", ""),
    }

    # 진도 좌우 균형 분배: 좌측=ceil(n/2), 우측=나머지. 좌측 위→아래가 앞 회차, 그다음 우측.
    progress_rows = lecture.get("progress", [])
    left_count, right_count = progress_split(len(progress_rows))
    for idx in range(1, 6):
        row = progress_rows[idx - 1] if idx <= left_count else {}
        placeholder_map[f"진도_좌_{idx}_날짜"] = progress_date_display(row) if row else ""
        placeholder_map[f"진도_좌_{idx}_내용"] = progress_content_display(row) if row else ""

    for idx in range(1, 6):
        source_index = left_count + idx - 1
        row = progress_rows[source_index] if idx <= right_count else {}
        placeholder_map[f"진도_우_{idx}_날짜"] = progress_date_display(row) if row else ""
        placeholder_map[f"진도_우_{idx}_내용"] = progress_content_display(row) if row else ""

    return {key: text_value(value) for key, value in placeholder_map.items()}


def replace_placeholders_in_text(text, placeholder_map):
    def replace(match):
        key = match.group(1).strip()
        return placeholder_map.get(key, "")

    return PLACEHOLDER_RE.sub(replace, text)


def replace_placeholders_in_paragraph(paragraph, placeholder_map):
    if not paragraph.runs:
        return
    original = "".join(run.text for run in paragraph.runs)
    if "{{" not in original:
        return
    replaced = replace_placeholders_in_text(original, placeholder_map)
    paragraph.runs[0].text = replaced
    for run in paragraph.runs[1:]:
        run.text = ""
    # 진도계획 휴강은 빨간색으로(휴강일 필드는 강의개요에서 빠져 여기 외엔 '휴강'으로 시작하는 셀 없음).
    if replaced.strip().startswith("휴강"):
        paragraph.runs[0].font.color.rgb = HOLIDAY_RED


def replace_placeholders_in_text_frame(text_frame, placeholder_map):
    for paragraph in text_frame.paragraphs:
        replace_placeholders_in_paragraph(paragraph, placeholder_map)


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def replace_placeholders_in_shape(shape, placeholder_map):
    if getattr(shape, "has_text_frame", False):
        replace_placeholders_in_text_frame(shape.text_frame, placeholder_map)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                replace_placeholders_in_text_frame(cell.text_frame, placeholder_map)


def replace_placeholders_in_slide(slide, placeholder_map):
    for shape in iter_shapes(slide.shapes):
        replace_placeholders_in_shape(shape, placeholder_map)


def collect_texts_from_slide(slide):
    texts = []
    for shape in iter_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            texts.append(shape.text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    texts.append(cell.text)
    return texts


def unresolved_placeholders(slide):
    found = []
    for text in collect_texts_from_slide(slide):
        found.extend(match.group(0) for match in PLACEHOLDER_RE.finditer(text or ""))
    return sorted(set(found))


def validate_required_values(lecture):
    reports = []
    fields = lecture.get("fields", {})
    required = {
        "메인 제목": fields.get("메인 제목", ""),
        "과목": fields.get("과목", ""),
        "강사명": fields.get("강사명", ""),
        "강좌명": fields.get("강좌명", ""),
    }
    for field, value in required.items():
        if not text_value(value).strip():
            reports.append(
                report_row(
                    "확인필요",
                    lecture,
                    field,
                    "REQUIRED_FIELD_EMPTY",
                    f"{field} 필수값이 비어 있습니다.",
                    suggestion="강좌 입력 시트에서 값을 입력해 주세요.",
                )
            )
    if not lecture.get("progress"):
        reports.append(
            report_row(
                "확인필요",
                lecture,
                "진도표",
                "REQUIRED_PROGRESS_EMPTY",
                "진도표가 비어 있습니다.",
                suggestion="최소 1개 이상의 진도 행을 입력해 주세요.",
            )
        )
    return reports


def validate_progress_overflow(lecture):
    if len(lecture.get("progress", [])) <= 10:
        return []
    return [
        report_row(
            "경고",
            lecture,
            "진도표",
            "PROGRESS_OVERFLOW",
            "진도표가 10개를 초과하여 PPTX에는 10개까지만 반영했습니다.",
            computed_value=str(len(lecture.get("progress", []))),
            suggestion="PPTX에서 추가 진도를 수동 반영하거나 템플릿 행을 확장해 주세요.",
        )
    ]


def generate_pptx_from_template(
    lectures,
    template_path,
    output_path,
    teacher_photo_dir=None,
):
    prs = Presentation(template_path)
    if not prs.slides:
        raise ValueError("Template PPTX has no slides.")

    template_slide = prs.slides[0]
    # 세로 중앙 정렬을 적용할 placeholder 값 박스 id(치환 전 템플릿 기준; 복제해도 id 유지).
    content_ids = {
        shape.shape_id
        for shape in template_slide.shapes
        if getattr(shape, "has_text_frame", False) and "{{" in shape.text_frame.text
    }
    reports = []
    generated_slides = []
    slide_subjects = []
    # 슬라이드는 프레젠테이션당 한 크기만 가능 → A4(템플릿 높이) 고정. 넘치는 슬라이드는 본문 폰트를 줄여 맞춘다.
    target_bottom = prs.slide_height - A4_BOTTOM_MARGIN

    for lecture in lectures:
        reports.extend(validate_required_values(lecture))
        reports.extend(validate_progress_overflow(lecture))
        slide = clone_template_slide(prs, template_slide)
        replace_placeholders_in_slide(slide, build_placeholder_map(lecture))
        left_count, right_count = progress_split(len(lecture.get("progress", [])))
        balance_progress_columns(slide, left_count, right_count)   # 안 쓰는 칸 도형 삭제
        fit_oneline_badges(slide)
        fit_title(slide)
        fit_scale, fitted = fit_slide_to_height(slide, content_ids, target_bottom)
        if not fitted:
            reports.append(
                report_row(
                    "확인필요",
                    lecture,
                    "PPTX",
                    "A4_OVERFLOW",
                    "강의특징/관리 본문 폰트를 가독성 하한까지 줄여도 A4 한 장에 안 들어갑니다. "
                    "코드가 더 줄이지 않으니(가독성 우선), 강사 분량을 줄여야 합니다.",
                    raw_value=f"scale={fit_scale}, prose_floor_cfg={PROSE_MIN_PT:.0f}pt",
                    suggestion="강의특징/관리 프로그램 문구를 줄여 주세요(이번 한 장 분량 초과).",
                )
            )
        teacher_name = lecture.get("fields", {}).get("강사명", "")
        inserted_photo = apply_teacher_photo(slide, teacher_name, teacher_photo_dir)
        if teacher_photo_dir and not inserted_photo:
            reports.append(
                report_row(
                    "확인필요",
                    lecture,
                    "강사 사진",
                    "TEACHER_PHOTO_NOT_FOUND",
                    f"'{teacher_name}' 강사 사진을 찾지 못해 회색 박스로 둡니다.",
                    suggestion=f"{teacher_photo_dir}/{teacher_name}.jpg 형태로 파일을 넣어 주세요.",
                )
            )
        remaining = unresolved_placeholders(slide)
        for placeholder in remaining:
            reports.append(
                report_row(
                    "경고",
                    lecture,
                    "PPTX",
                    "UNRESOLVED_PLACEHOLDER",
                    f"치환되지 않은 placeholder가 남아 있습니다: {placeholder}",
                    raw_value=placeholder,
                    suggestion="placeholder_map 또는 템플릿 placeholder명을 확인해 주세요.",
                )
            )
        generated_slides.append(slide)
        # 배지와 동일하게 정규화한 과목으로 색 테두리 매칭(미적분·확통·기하 등 세부과목 → '수학').
        slide_subjects.append(normalize_subject(lecture.get("fields", {}).get("과목", "")))

    remove_slide(prs, template_slide)
    # 슬라이드 높이는 A4(템플릿 그대로) 유지 — fit_slide_to_height가 내용을 A4 안에 맞춰 둠.
    # 최종 슬라이드 크기가 정해진 뒤 과목 색 테두리를 두른다(테두리가 전체 면적을 감싸야 하므로).
    for slide, subject in zip(generated_slides, slide_subjects):
        add_subject_band(slide, subject, prs.slide_width, prs.slide_height)
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return Path(output_path), reports


def generate_pptx(lectures, output_dir, template_path=None, teacher_photo_dir=None):
    output = Path(output_dir) / "generated_pptx"
    output.mkdir(parents=True, exist_ok=True)
    template = Path(template_path or DEFAULT_TEMPLATE_PATH)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M")
    output_path = output / f"강의계획서_초안_{timestamp}.pptx"
    return generate_pptx_from_template(
        lectures=lectures,
        template_path=template,
        output_path=output_path,
        teacher_photo_dir=teacher_photo_dir,
    )
