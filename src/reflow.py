# -*- coding: utf-8 -*-
"""글자 수에 맞춰 박스(및 같은 행의 배경/테두리)를 자동 확장하고 아래 요소를 밀어내는 리플로우.

전략(B안):
- 슬라이드의 도형을 세로로 겹치는 것끼리 '행(band)'으로 묶는다.
- 각 행에서 글자 박스가 필요로 하는 높이를 추정해 행을 키운다.
- 행을 채우는 도형(배경 셀·테두리·글자 박스)은 함께 키우고, 위 행이 커진 만큼 아래 행을 내린다.
- 슬라이드 높이는 전체 증가분만큼 늘린다.

줄 수는 실제 렌더링 없이 글자 폭으로 추정한다(근사치). 한글 1글자=폰트크기, 영문/숫자=절반.
"""
import math

from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu

try:
    from src.teacher_photo import is_photo_box
except ModuleNotFoundError:
    from .teacher_photo import is_photo_box

EMU_PER_PT = 12700.0
EMU_PER_CM = 360000.0

LINE_SLOT = 1.22        # 한 줄이 차지하는 높이(폰트 pt 배수, 실제 렌더 줄높이에 근접). 과대 추정 시 박스가 떠 보임.
DEFAULT_FONT_PT = 9.0
PAD = 0.06 * EMU_PER_CM            # 박스 안 글자 상·하 여백(각 측). 중앙 정렬 시 상하 여백 일정(과대 시 정보 그리드까지 부풂).
OVERFLOW_TOL = 0.10 * EMU_PER_CM   # 이만큼 넘쳐야 확장(미세 넘침 무시)
MIN_GROW = 0.06 * EMU_PER_CM       # 이보다 작은 확장은 0으로 스냅
THIN_HEIGHT = 0.30 * EMU_PER_CM    # 이보다 얇은 도형은 장식선 → 키우지 않음
REACH_TOL = 0.18 * EMU_PER_CM      # 행 바닥에 이만큼 근접하면 '행을 채우는' 도형


def _first_font_pt(shape):
    if not shape.has_text_frame:
        return DEFAULT_FONT_PT
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                return run.font.size.pt
    for para in shape.text_frame.paragraphs:
        if para.font.size:
            return para.font.size.pt
    return DEFAULT_FONT_PT


def _insets_pt(text_frame):
    def val(margin, default):
        return (margin if margin is not None else default) / EMU_PER_PT
    return (
        val(text_frame.margin_top, Emu(45720)),
        val(text_frame.margin_bottom, Emu(45720)),
        val(text_frame.margin_left, Emu(91440)),
        val(text_frame.margin_right, Emu(91440)),
    )


def _est_lines(text, width_pt, font_pt, wrap=True):
    text = (text or "").rstrip("\n \t")
    if not text:
        return 0
    if not wrap:
        # 자동 줄바꿈 off → 명시적 줄바꿈(\n) 수만큼만 줄로 센다(가로 넘침은 박스 높이와 무관).
        return max(1, len(text.split("\n")))
    total = 0
    for line in text.split("\n"):
        if not line:
            total += 1
            continue
        visual = sum((font_pt * 0.5) if ord(ch) < 128 else (font_pt * 1.0) for ch in line)
        total += max(1, math.ceil(visual / max(1.0, width_pt)))
    return max(1, total)


def _needed_bottom_emu(shape):
    """글자 박스가 내용을 모두 담으려면 도형 바닥이 어디까지 내려가야 하는지(EMU)."""
    tf = shape.text_frame
    if not tf.text.strip():
        return shape.top + shape.height
    mt, mb, ml, mr = _insets_pt(tf)
    font_pt = _first_font_pt(shape)
    width_pt = shape.width / EMU_PER_PT - ml - mr
    wrap = tf.word_wrap if tf.word_wrap is not None else True
    lines = _est_lines(tf.text, width_pt, font_pt, wrap)
    text_h_pt = lines * font_pt * LINE_SLOT
    needed_h_emu = (mt + text_h_pt + mb) * EMU_PER_PT + 2 * PAD
    return shape.top + int(round(needed_h_emu))


def _cluster_bands(shapes):
    """세로로 겹치는 도형끼리 행으로 묶는다. 반환: [(top, bottom, [shapes])] (top 오름차순)."""
    items = sorted(shapes, key=lambda s: s.top)
    bands = []
    for sh in items:
        top, bottom = sh.top, sh.top + sh.height
        if bands and top < bands[-1][1]:   # 직전 행 바닥보다 위에서 시작 → 겹침
            b = bands[-1]
            b[0] = min(b[0], top)
            b[1] = max(b[1], bottom)
            b[2].append(sh)
        else:
            bands.append([top, bottom, [sh]])
    return bands


def reflow_slide(slide, content_ids=None):
    """슬라이드를 제자리 리플로우. 늘어난 총 높이(EMU)를 반환.

    content_ids: 세로 중앙 정렬을 적용할 도형 id 집합(placeholder 값 박스). None이면 정렬은 건드리지 않음.
    """
    shapes = [sh for sh in slide.shapes if sh.top is not None and sh.height is not None]
    if not shapes:
        return 0

    # placeholder 값 박스는 세로 중앙 정렬로 통일(글자가 박스 위/아래 가장자리에 붙지 않도록).
    if content_ids:
        for sh in shapes:
            if sh.shape_id in content_ids and sh.has_text_frame:
                sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    orig = {id(sh): (sh.top, sh.height) for sh in shapes}
    bands = _cluster_bands(shapes)

    # 1) 각 행의 확장량 계산
    deltas = []
    for top, bottom, members in bands:
        needed_bottom = bottom
        for sh in members:
            if sh.has_text_frame and sh.text_frame.text.strip():
                nb = _needed_bottom_emu(sh)
                if nb > needed_bottom:
                    needed_bottom = nb
        delta = needed_bottom - bottom
        if delta < OVERFLOW_TOL or delta < MIN_GROW:
            delta = 0
        deltas.append(delta)

    # 2) 누적 오프셋으로 아래 행을 밀고, 행을 채우는 도형은 함께 확장
    offset = 0
    for (top, bottom, members), delta in zip(bands, deltas):
        band_mid = (top + bottom) / 2.0
        for sh in members:
            otop, oh = orig[id(sh)]
            obottom = otop + oh
            is_thin = oh < THIN_HEIGHT
            # 얇으면서 행의 아래쪽에 붙은 도형만 '바닥 장식선' → 행이 커진 만큼 추가로 내림.
            bottom_decoration = is_thin and otop > band_mid
            if bottom_decoration:
                sh.top = int(round(otop + offset + delta))
            else:
                sh.top = int(round(otop + offset))
            # 행을 세로로 채우는 도형(배경 셀·테두리·글자 박스)만 확장.
            # 얇은 장식선과 우상단 사진 박스(정사각 유지 필요)는 제외.
            if delta and not is_thin and obottom >= bottom - REACH_TOL and not is_photo_box(sh):
                sh.height = int(round(oh + delta))
        offset += delta

    return offset
