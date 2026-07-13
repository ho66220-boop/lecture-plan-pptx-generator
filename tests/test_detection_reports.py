# -*- coding: utf-8 -*-
"""Batch 3 — 탐지망 복구 회귀 테스트 (P3-2 청구 미확정 · P2-3 미매핑 placeholder · P3-5 제목 오버플로).

공통 주제: "조용히 넘어가는 것에 리포트 달기". 세 건 모두 값·치환·폰트 동작은
불변이고 validation_report 행만 추가된다 — 경계 테스트(T3·T5)로 오탐 부재를 잠근다.
"""
from pptx import Presentation
from pptx.util import Emu, Pt

try:
    from src.generate_pptx import generate_pptx_from_template
    from src.normalize import normalize_lecture
except ModuleNotFoundError:
    from ..src.generate_pptx import generate_pptx_from_template
    from ..src.normalize import normalize_lecture

CM = 360000  # EMU


def make_lecture_raw(gubun="정규반", season=""):
    fields = {"강사명": "홍길동", "강좌명": "테스트강좌", "구분": gubun,
              "강의형태": "현장강의", "과목": "국어", "수업 요일 / 시간": "화 19:00"}
    if season:
        fields["시즌"] = season
    return {"source_file": "t.xlsx", "source_sheet": "강좌1", "fields": fields,
            "progress": [{"날짜": d, "수업 주제": "x"} for d in ("7/7", "7/14")]}


def build_template(path, texts, title_pt=None):
    """텍스트박스들로 최소 템플릿 pptx 생성. title_pt를 주면 마지막 박스를
    큰 제목 규격(폭 11cm ≥ TITLE_MIN_WIDTH, 상단, 지정 pt)으로 만든다."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i, text in enumerate(texts):
        is_title = title_pt is not None and i == len(texts) - 1
        width = Emu(11 * CM) if is_title else Emu(6 * CM)
        tb = slide.shapes.add_textbox(Emu(1 * CM), Emu(int((1 + i * 1.2) * CM)), width, Emu(int(0.9 * CM)))
        tb.text_frame.text = text
        if is_title:
            for run in tb.text_frame.paragraphs[0].runs:
                run.font.size = Pt(title_pt)
    prs.save(path)
    return path


def run_template(tmp_path, texts, lecture=None, title_pt=None):
    template = build_template(tmp_path / "template.pptx", texts, title_pt=title_pt)
    lec, _ = normalize_lecture(lecture or make_lecture_raw(), 1, 2026)
    _, reports = generate_pptx_from_template(
        [lec], template, tmp_path / "out.pptx", teacher_photo_dir=None
    )
    return lec, reports


# ══ T1 — P3-2: 구분·시즌 공백 → 청구 미확정 리포트 (계약 불변 동시 단언) ══

def test_billing_undetermined_reported_when_no_basis():
    """# P3-2 회귀 방지
    구분·시즌이 모두 비면 monthly로 조용히 가정된다(특강 표기 누락 = 오청구 위험).
    기대: BILLING_UNDETERMINED 리포트 존재 + billing 반환은 여전히 monthly
    (classify_billing 17개 테스트 계약 불변)."""
    lec, reports = normalize_lecture(make_lecture_raw(gubun="", season=""), 1, 2026)

    codes = [r["issue_code"] for r in reports]
    assert "BILLING_UNDETERMINED" in codes                 # 침묵 가정 금지
    assert lec["billing"] == "monthly"                     # 계산·계약은 불변


# ══ T2 — P2-3: 템플릿의 미매핑 placeholder가 리포트로 드러난다 ══

def test_unmapped_placeholder_reported(tmp_path):
    """# P2-3 회귀 방지
    Phase 2 실증: {{존재하지않는키}}는 조용히 ""로 치환되고 기존 unresolved 검출은
    빈 리스트만 반환(죽은 코드). 기대: 치환 전 대조로 미매핑 키가 리포트에 남는다.
    치환 동작 자체는 불변('홍길동 / ' 렌더 유지)도 함께 단언."""
    _, reports = run_template(tmp_path, ["{{강사명}} / {{존재하지않는키}}"])

    unmapped = [r for r in reports if r["issue_code"] == "UNMAPPED_PLACEHOLDER"]
    assert len(unmapped) == 1                              # 실행당 키당 1회
    assert "존재하지않는키" in unmapped[0]["raw_value"]
    # 치환 결과는 기존과 동일해야 한다(동작 변경 아님 — 빈칸 렌더 유지).
    out = Presentation(tmp_path / "out.pptx")
    texts = " ".join(sh.text_frame.text for sh in out.slides[0].shapes
                     if getattr(sh, "has_text_frame", False))
    assert "홍길동 /" in texts and "존재하지않는키" not in texts


# ══ T3 — P2-3 경계: 키는 있고 값이 빈 placeholder는 오탐하지 않는다 ══

def test_mapped_but_empty_value_not_reported(tmp_path):
    """# P2-3 경계 잠금
    {{휴강일}}은 map에 키가 있고 값이 ""(입력 안 함)일 뿐 — 의도된 빈칸이다.
    미매핑과 혼동해 리포트하면 결측 필드마다 스팸이 된다."""
    _, reports = run_template(tmp_path, ["휴강: {{휴강일}}"])

    codes = [r["issue_code"] for r in reports]
    assert "UNMAPPED_PLACEHOLDER" not in codes


# ══ T4 — P3-5: 폰트 하한에서도 줄 수를 초과하는 제목 → 리포트 ══

def test_title_overflow_reported_at_font_floor(tmp_path):
    """# P3-5 회귀 방지
    초장문 강좌명은 fit_title이 14pt 하한까지만 줄이고, 그래도 2줄을 넘치면
    아래 요소와 겹칠 수 있는데 어떤 리포트도 없었다(A4_OVERFLOW는 하단 기준이라
    상단 제목을 못 잡음). 기대: TITLE_OVERFLOW(확인필요) 리포트 — 'COM 검수 시
    이 슬라이드 먼저 보라'는 신호."""
    long_title_raw = make_lecture_raw()
    long_title_raw["fields"]["강좌명"] = "가나다라마바사아자차" * 6   # 60자
    _, reports = run_template(tmp_path, ["{{메인제목}}"], lecture=long_title_raw, title_pt=20)

    codes = [r["issue_code"] for r in reports]
    assert "TITLE_OVERFLOW" in codes


# ══ T5 — P3-5 경계: 통상 길이 제목은 미발동 ══

def test_normal_title_not_reported(tmp_path):
    """# P3-5 경계 잠금 — 통상 제목(2줄 안에 들어감)에 오탐하지 않는다."""
    _, reports = run_template(tmp_path, ["{{메인제목}}"], title_pt=20)

    codes = [r["issue_code"] for r in reports]
    assert "TITLE_OVERFLOW" not in codes
