# -*- coding: utf-8 -*-
"""Batch 4 — 저위험 정리 검증 (P3-4 파일명 초단위 · P2-9 사진 메시지 보강).

나머지 항목(P2-7 중복 통합·P2-8 죽은 라벨 삭제·P2-11 레이아웃 가드)은
동작 불변 정리라 기존 116개 스위트가 안전망이다.
"""
import re

from pptx import Presentation
from pptx.util import Emu

try:
    from src.generate_pptx import generate_pptx, generate_pptx_from_template
    from src.normalize import normalize_lecture
except ModuleNotFoundError:
    from ..src.generate_pptx import generate_pptx, generate_pptx_from_template
    from ..src.normalize import normalize_lecture

CM = 360000


def make_lecture():
    raw = {"source_file": "t.xlsx", "source_sheet": "강좌1",
           "fields": {"강사명": "김테스트", "강좌명": "테스트강좌", "구분": "정규반",
                      "강의형태": "현장강의", "과목": "국어", "수업 요일 / 시간": "화 19:00"},
           "progress": [{"날짜": "7/7", "수업 주제": "x"}]}
    lec, _ = normalize_lecture(raw, 1, 2026)
    return lec


def build_template(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Emu(CM), Emu(CM), Emu(6 * CM), Emu(CM))
    tb.text_frame.text = "{{강사명}}"
    prs.save(path)
    return path


def test_pptx_filename_has_seconds(tmp_path):
    """# P3-4 회귀 방지
    같은 분(minute) 안에 재실행하면 파일명이 겹쳐 조용히 덮어쓰던 문제 —
    타임스탬프는 초 단위(YYYYMMDD_HHMMSS)여야 한다. 고정 이름 산출물
    (normalized_data/validation_report)은 의도적으로 그대로 둔다."""
    template = build_template(tmp_path / "template.pptx")
    pptx_path, _ = generate_pptx([make_lecture()], tmp_path, template_path=template)

    assert re.fullmatch(r"강의계획서_초안_\d{8}_\d{6}\.pptx", pptx_path.name), pptx_path.name


def test_photo_not_found_message_mentions_box_possibility(tmp_path):
    """# P2-9 회귀 방지
    사진 미삽입은 '파일 없음'뿐 아니라 '템플릿 사진박스가 기준 위치를 벗어나
    미식별'일 수도 있다 — NOT_FOUND 메시지가 그 가능성을 함께 안내해야
    오탐 메시지(파일이 있는데 없다고 단정)가 되지 않는다. 발동 조건은 불변."""
    template = build_template(tmp_path / "template.pptx")
    photo_dir = tmp_path / "photos"
    photo_dir.mkdir()
    _, reports = generate_pptx_from_template(
        [make_lecture()], template, tmp_path / "out.pptx", teacher_photo_dir=str(photo_dir)
    )

    not_found = [r for r in reports if r["issue_code"] == "TEACHER_PHOTO_NOT_FOUND"]
    assert len(not_found) == 1                              # 발동 조건 불변
    combined = not_found[0]["message"] + not_found[0]["suggestion"]
    assert "사진박스" in combined                            # 미식별 가능성 안내 포함
