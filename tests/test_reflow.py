# -*- coding: utf-8 -*-
"""reflow_slide 견고성 테스트 — width=None 도형이 섞여도 크래시하지 않는다.

reflow는 _needed_bottom_emu에서 shape.width로 나눗셈을 하므로, 좌표가 일부만
정의된(width=None) 도형이 필터를 통과하면 TypeError가 난다. 가드가 그런 도형을
처리 전에 걸러내는지 확인한다.
"""
from pptx import Presentation
from pptx.util import Emu

try:
    from src.reflow import reflow_slide
except ModuleNotFoundError:
    from ..src.reflow import reflow_slide


class _WidthNoneShape:
    """top·height는 정상이지만 width만 None인 도형(좌표 일부 미상).
    가드가 작동하면 필터에서 빠져 아래 has_text_frame은 절대 호출되지 않는다."""
    top = 1_500_000
    height = 500_000
    width = None

    @property
    def has_text_frame(self):  # pragma: no cover - 호출되면 가드 실패
        raise AssertionError("width=None 도형이 필터를 통과해 처리됨 — 가드 미작동")


class _FakeSlide:
    def __init__(self, shapes):
        self.shapes = shapes


def _new_slide_with_textbox():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    tb = slide.shapes.add_textbox(Emu(200_000), Emu(200_000), Emu(3_000_000), Emu(600_000))
    tb.text_frame.text = "정상 텍스트 박스"
    return slide


def test_width_none_shape_is_skipped_without_crash():
    """정상 텍스트 박스 + width=None 도형을 함께 넣어도 예외 없이 처리되고,
    width=None 도형은 처리 대상에서 제외된다(has_text_frame 미호출)."""
    slide = _new_slide_with_textbox()
    fake = _FakeSlide(list(slide.shapes) + [_WidthNoneShape()])
    grew = reflow_slide(fake, content_ids=None)   # 가드 없으면 여기서 AssertionError/TypeError
    assert isinstance(grew, int)
    assert grew >= 0


def test_reflow_runs_on_normal_slide():
    """정상 슬라이드에서 reflow가 예외 없이 정수(증가 높이)를 돌려준다."""
    slide = _new_slide_with_textbox()
    grew = reflow_slide(slide, content_ids=None)
    assert isinstance(grew, int)
    assert grew >= 0
