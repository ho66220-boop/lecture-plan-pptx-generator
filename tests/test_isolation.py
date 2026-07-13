# -*- coding: utf-8 -*-
"""Batch 1 — 격리·리포트 보존 회귀 테스트 (P2-2 · P5-2 · P3-3).

원칙: 강좌 1건의 예외가 배치 전체를 죽이지 않고(격리), 어떤 예외 경로에서도
validation_report가 파일로 남으며(보존), 산출물 저장 실패는 스택트레이스가 아니라
명확한 안내 또는 대체명 저장으로 처리된다.

주의: 감사 프롬프트의 T1(2/29→LECTURE_NORMALIZE_FAILED)은 P2-2 수정이 2층
(행 단위 replace 가드 + 강좌 단위 격리)이라 가드가 먼저 잡으면 강좌가 살아남는
더 나은 결과가 된다. 그래서 T1을 T1a(행 가드)/T1b(범용 격리)로 분리했다.
"""
import os

import pytest
from openpyxl import Workbook, load_workbook

try:
    from src import generate_pptx as gp_mod
    from src import normalize as normalize_mod
    from src.normalize import normalize_lectures
    from src.pipeline import run_pipeline
except ModuleNotFoundError:
    from ..src import generate_pptx as gp_mod
    from ..src import normalize as normalize_mod
    from ..src.normalize import normalize_lectures
    from ..src.pipeline import run_pipeline


# ── 합성 입력 헬퍼 ──

def make_raw(dates, teacher="김강사", course="테스트강좌", sheet="강좌1"):
    return {
        "source_file": "t.xlsx",
        "source_sheet": sheet,
        "fields": {
            "강사명": teacher, "강좌명": course, "구분": "정규반",
            "강의형태": "현장강의", "과목": "국어", "수업 요일 / 시간": "화 19:00",
        },
        "progress": [{"날짜": d, "수업 주제": f"{i + 1}강"} for i, d in enumerate(dates)],
    }


def build_input_xlsx(path, courses):
    """courses: [(시트명, 강좌명, 진도날짜들)] → 강좌 시트만 있는 입력 엑셀 생성."""
    wb = Workbook()
    wb.remove(wb.active)
    for sheet, course, dates in courses:
        ws = wb.create_sheet(sheet)
        for label, value in (
            ("강사명", "김강사"), ("강좌명", course), ("메인 제목", course),
            ("과목", "국어"), ("구분", "정규반"), ("강의형태", "현장강의"),
            ("수업 요일 / 시간", "화 19:00"),
        ):
            ws.append([label, value])
        ws.append(["회차", "날짜", "수업 주제", "상세 내용", "비고"])
        for i, d in enumerate(dates, start=1):
            ws.append([f"{i}회", d, f"{i}강", "", ""])
    wb.save(path)
    return path


def read_issue_codes(validation_xlsx):
    ws = load_workbook(validation_xlsx).active
    rows = list(ws.iter_rows(values_only=True))
    idx = {h: i for i, h in enumerate(rows[0])}
    return [r[idx["issue_code"]] for r in rows[1:]]


# ══ T1a — P2-2 행 가드: 윤년 2/29 롤오버가 크래시 대신 그 행만 실패 처리 ══

def test_leap_rollover_row_guarded_batch_survives():
    """# P2-2 회귀 방지 (행 단위 replace 가드)
    base_year 윤년(2028) + 진도 12/20→1/10→2/29: 12→1 wrap으로 연도가 +1된 뒤
    2/29가 평년(2029)으로 보정되며 replace에서 ValueError가 난다(가드 없으면
    배치 전체가 죽음). 기대: 예외 없이 두 강좌 모두 산출되고, 2/29 행만
    DATE_PARSE_FAILED로 리포트된다(강좌는 생존).
    (Batch 2에서 롤오버 판별이 '12→1 인접 wrap만'으로 좁혀져, 구 픽스처
    12/20→2/29는 wrap이 아닌 OUT_OF_ORDER가 되므로 wrap 경유 경로로 갱신함.)"""
    crash = make_raw(["12/20", "1/10", "2/29"], course="윤년크래시", sheet="강좌1")
    normal = make_raw(["12/21", "12/28"], course="정상강좌", sheet="강좌2")
    lectures, reports = normalize_lectures([crash, normal], 2028)   # 예외 전파 시 여기서 실패

    names = [lec["fields"]["강좌명"] for lec in lectures]
    assert "정상강좌" in names                       # 정상 강좌 생존
    assert "윤년크래시" in names                     # 크래시 강좌도 행 가드로 생존
    codes = [r["issue_code"] for r in reports]
    assert "DATE_PARSE_FAILED" in codes              # 2/29 행이 실패로 리포트됨


# ══ T1b — P2-2 강좌 격리: 정규화 내부 예외가 배치를 죽이지 않는다 ══

def test_normalize_lecture_exception_isolated(monkeypatch):
    """# P2-2 회귀 방지 (강좌 단위 try/except)
    특정 강좌에서만 내부 함수(calculate_fee)가 예외를 던지게 하면, 현재는 배치
    전체가 죽는다. 기대: 그 강좌만 LECTURE_NORMALIZE_FAILED로 스킵되고
    정상 강좌는 산출되며 예외가 호출자로 전파되지 않는다."""
    real_fee = normalize_mod.calculate_fee

    def bomb_fee(lecture_id, *args, **kwargs):
        if kwargs.get("teacher_name") == "폭탄강사":
            raise RuntimeError("합성 예외(테스트)")
        return real_fee(lecture_id, *args, **kwargs)

    monkeypatch.setattr(normalize_mod, "calculate_fee", bomb_fee)

    bomb = make_raw(["7/7", "7/14"], teacher="폭탄강사", course="폭탄강좌", sheet="강좌1")
    normal = make_raw(["7/8", "7/15"], teacher="정상강사", course="정상강좌", sheet="강좌2")
    lectures, reports = normalize_lectures([bomb, normal], 2026)   # 예외 전파 시 여기서 실패

    names = [lec["fields"]["강좌명"] for lec in lectures]
    assert "정상강좌" in names                       # 정상 강좌 생존
    assert "폭탄강좌" not in names                   # 실패 강좌는 스킵
    failed = [r for r in reports if r["issue_code"] == "LECTURE_NORMALIZE_FAILED"]
    assert len(failed) == 1                          # 격리 리포트 존재
    assert failed[0]["source_sheet"] == "강좌1"      # 누가 실패했는지 특정 가능
    assert "RuntimeError" in failed[0]["message"]    # 왜 실패했는지(예외 타입) 포함


# ══ T2 — P5-2: PPTX 루프 강좌 격리 + validation_report 보존 ══

def test_pptx_loop_isolates_lecture_and_preserves_report(tmp_path, monkeypatch):
    """# P5-2 회귀 방지
    PPTX 생성 중 강좌 1건에서만 예외(fit 단계)가 나도록 하면, 현재는 run 전체가
    죽고 validation_report.xlsx도 안 남는다. 기대: 정상 강좌 2건의 슬라이드가
    생성되고, PPTX_GENERATION_FAILED 리포트가 남으며, run_pipeline이 예외 없이
    완료되고 validation_report.xlsx가 파일로 존재한다."""
    input_path = build_input_xlsx(
        tmp_path / "input.xlsx",
        [("강좌1", "정상A", ["7/7", "7/14"]),
         ("강좌2", "폭탄강좌", ["7/8", "7/15"]),
         ("강좌3", "정상B", ["7/9", "7/16"])],
    )
    real_fit = gp_mod.fit_slide_to_height

    def bomb_fit(slide, content_ids, target_bottom, **kwargs):
        texts = " ".join(
            sh.text_frame.text for sh in slide.shapes
            if getattr(sh, "has_text_frame", False)
        )
        if "폭탄강좌" in texts:
            raise RuntimeError("합성 fit 예외(테스트)")
        return real_fit(slide, content_ids, target_bottom, **kwargs)

    monkeypatch.setattr(gp_mod, "fit_slide_to_height", bomb_fit)

    result = run_pipeline(input_path=str(input_path), output_dir=str(tmp_path / "out"),
                          base_year=2026)            # 예외 전파 시 여기서 실패

    assert os.path.exists(result["validation_report"])            # 리포트 파일 보존
    codes = read_issue_codes(result["validation_report"])
    assert "PPTX_GENERATION_FAILED" in codes                      # 격리 리포트 존재
    assert result["pptx_path"] and os.path.exists(result["pptx_path"])
    from pptx import Presentation
    assert len(Presentation(result["pptx_path"]).slides) == 2     # 정상 2강좌만 슬라이드


# ══ T3 — P3-3: 산출물 저장 실패의 친절 처리 ══

def _fail_primary_save(monkeypatch, primary_basename):
    """Workbook.save가 지정 파일명(기본 산출명)에만 PermissionError를 던지게 패치."""
    real_save = Workbook.save

    def guarded_save(self, filename):
        if os.path.basename(str(filename)) == primary_basename:
            raise PermissionError(13, "파일이 다른 프로세스에서 사용 중(합성)")
        return real_save(self, filename)

    monkeypatch.setattr(Workbook, "save", guarded_save)


def test_normalized_xlsx_locked_falls_back_to_alt_name(tmp_path, monkeypatch):
    """# P3-3 회귀 방지 (openpyxl 저장 경로 ①)
    normalized_data.xlsx가 잠겨 있으면 현재는 PermissionError 스택트레이스로 중단.
    기대: 대체 파일명으로 저장하고 run이 완료되며 결과 경로가 실제 존재한다."""
    input_path = build_input_xlsx(tmp_path / "input.xlsx", [("강좌1", "정상A", ["7/7"])])
    _fail_primary_save(monkeypatch, "normalized_data.xlsx")

    result = run_pipeline(input_path=str(input_path), output_dir=str(tmp_path / "out"),
                          base_year=2026, make_pptx=False)        # PermissionError 전파 시 실패

    assert os.path.exists(result["normalized_xlsx"])              # 대체명 파일이 실존
    assert os.path.basename(result["normalized_xlsx"]) != "normalized_data.xlsx"
    assert os.path.exists(result["validation_report"])            # 리포트는 정상 저장


def test_validation_report_locked_falls_back_to_alt_name(tmp_path, monkeypatch):
    """# P3-3 회귀 방지 (openpyxl 저장 경로 ② — 최후의 보루)
    validation_report.xlsx가 잠겨 있어도 대체 파일명 폴백이 실제로 파일을 남긴다."""
    input_path = build_input_xlsx(tmp_path / "input.xlsx", [("강좌1", "정상A", ["7/7"])])
    _fail_primary_save(monkeypatch, "validation_report.xlsx")

    result = run_pipeline(input_path=str(input_path), output_dir=str(tmp_path / "out"),
                          base_year=2026, make_pptx=False)        # PermissionError 전파 시 실패

    assert os.path.exists(result["validation_report"])            # 폴백 파일이 실존
    assert os.path.basename(result["validation_report"]) != "validation_report.xlsx"


def test_pptx_save_locked_reports_and_run_completes(tmp_path, monkeypatch):
    """# P3-3 회귀 방지 (python-pptx 저장 경로)
    PPTX 저장이 계속 실패해도(파일 잠금 등) run은 완료되고, OUTPUT_SAVE_FAILED
    리포트가 남으며 validation_report.xlsx는 파일로 존재한다(P5-2 보존과 결합)."""
    import pptx.presentation

    input_path = build_input_xlsx(tmp_path / "input.xlsx", [("강좌1", "정상A", ["7/7"])])

    def always_fail_save(self, file):
        raise PermissionError(13, "파일이 다른 프로세스에서 사용 중(합성)")

    monkeypatch.setattr(pptx.presentation.Presentation, "save", always_fail_save)

    result = run_pipeline(input_path=str(input_path), output_dir=str(tmp_path / "out"),
                          base_year=2026)            # PermissionError 전파 시 여기서 실패

    assert result["pptx_path"] == ""                              # pptx는 산출 실패로 빈 값
    assert os.path.exists(result["validation_report"])            # 리포트 파일은 보존
    codes = read_issue_codes(result["validation_report"])
    assert "OUTPUT_SAVE_FAILED" in codes                          # 실패가 리포트로 남음
